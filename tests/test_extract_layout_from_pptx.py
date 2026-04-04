from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches


SCRIPT_DIR = Path(__file__).resolve().parents[1] / "scripts"
sys.path.insert(0, str(SCRIPT_DIR))

from extract_layout_from_pptx import build_manifest_from_pptx  # noqa: E402


class ExtractLayoutFromPptxTests(unittest.TestCase):
    def test_extracts_page_metrics_from_real_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            deck_path = Path(tmp) / "sample.pptx"
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(5.6), Inches(0.8)).text = "标题"
            slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.0), Inches(3.2), Inches(1.0)).text = "卡片 A"
            slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.4), Inches(3.2), Inches(1.0)).text = "卡片 B"
            slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.8), Inches(2.4), Inches(0.35), Inches(0.35))
            slide.shapes.add_connector(1, Inches(5.0), Inches(2.75), Inches(5.0), Inches(4.2))
            prs.save(deck_path)

            state = {"pages": [{"page_id": "slide_01", "role": "hero_cover"}]}
            manifest = build_manifest_from_pptx(deck_path, state, {"pages": []})
            page = manifest["pages"][0]

            self.assertEqual(page["page_id"], "slide_01")
            self.assertEqual(page["role"], "hero_cover")
            self.assertGreater(page["occupancy"]["ratio"], 0.1)
            self.assertGreaterEqual(len(page["cards"]), 2)


if __name__ == "__main__":
    unittest.main()
