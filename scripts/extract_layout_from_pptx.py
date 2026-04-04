#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROLE_TO_ARCHETYPE = {
    "hero_cover": "hero_cover",
    "hero_problem": "diagnostic_board",
    "hero_proof": "proof_board",
    "hero_system": "system_map",
    "hero_diff": "comparison_board",
    "hero_value": "value_stack",
    "hero_cta": "cta_board",
}


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8")) if path.exists() else {}


def save_json(path: Path, payload: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def shape_name(shape) -> str:
    return str(getattr(shape, "name", "") or "").strip()


def shape_text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return " ".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()).strip()
    return ""


def local_tag(shape) -> str:
    tag = getattr(shape.element, "tag", "")
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def is_group(shape) -> bool:
    return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP


def iter_shapes(shapes):
    for shape in shapes:
        if is_group(shape):
            try:
                yield from iter_shapes(shape.shapes)
            except Exception:
                yield shape
        else:
            yield shape


def bbox_norm(shape, slide_w: float, slide_h: float) -> dict:
    left = float(shape.left) / slide_w
    top = float(shape.top) / slide_h
    width = float(shape.width) / slide_w
    height = float(shape.height) / slide_h
    return {
        "x": left,
        "y": top,
        "w": width,
        "h": height,
        "cx": left + width / 2,
        "cy": top + height / 2,
        "area": width * height,
    }


def classify_line_like(item: dict) -> bool:
    tag = item.get("tag", "")
    if tag == "cxnSp":
        return True
    shape_type_int = item.get("shape_type_int")
    if shape_type_int is not None and shape_type_int == int(MSO_SHAPE_TYPE.LINE):
        return True
    box = item.get("box", {})
    return (box.get("w", 0) < 0.02 and box.get("h", 0) > 0.12) or (box.get("h", 0) < 0.02 and box.get("w", 0) > 0.12)


def classify_circle_like(item: dict) -> bool:
    if item.get("tag", "") != "sp":
        return False
    box = item.get("box", {})
    area = box.get("area", 0)
    if area < 0.001 or area > 0.03:
        return False
    w, h = box.get("w", 0), box.get("h", 0)
    return abs(w - h) <= min(w, h) * 0.25 and not item.get("text", "")


def classify_card_like(item: dict) -> bool:
    if classify_line_like(item):
        return False
    box = item.get("box", {})
    return bool(item.get("text", "")) and box.get("w", 0) >= 0.12 and box.get("h", 0) >= 0.05


def is_noise(box: dict) -> bool:
    return box["area"] < 0.0008 or (box["y"] > 0.92 and box["h"] < 0.05)


def cluster_by_coordinate(items: list[dict], key: str, tolerance: float, min_size: int = 3) -> list[list[dict]]:
    if not items:
        return []
    items = sorted(items, key=lambda item: item[key])
    clusters: list[list[dict]] = [[items[0]]]
    for item in items[1:]:
        if abs(item[key] - clusters[-1][-1][key]) <= tolerance:
            clusters[-1].append(item)
        else:
            clusters.append([item])
    return [cluster for cluster in clusters if len(cluster) >= min_size]


def infer_line_endpoints(box: dict) -> tuple[dict, dict]:
    if box["h"] > box["w"] * 3:
        return (
            {"x": box["cx"], "y": box["y"]},
            {"x": box["cx"], "y": box["y"] + box["h"]},
        )
    if box["w"] > box["h"] * 3:
        return (
            {"x": box["x"], "y": box["cy"]},
            {"x": box["x"] + box["w"], "y": box["cy"]},
        )
    return (
        {"x": box["x"], "y": box["y"]},
        {"x": box["x"] + box["w"], "y": box["y"] + box["h"]},
    )


def nearest_anchor(point: dict, shapes: list[dict]) -> dict:
    best = None
    best_dist = float("inf")
    for item in shapes:
        b = item["box"]
        anchor_x = min(max(point["x"], b["x"]), b["x"] + b["w"])
        anchor_y = min(max(point["y"], b["y"]), b["y"] + b["h"])
        dist = math.hypot(point["x"] - anchor_x, point["y"] - anchor_y)
        if dist < best_dist:
            best_dist = dist
            best = {"x": anchor_x, "y": anchor_y, "shape_id": item["id"]}
    return best or {"x": point["x"], "y": point["y"], "shape_id": ""}


def build_page_entry(slide, page_id: str, role: str, slide_w: float, slide_h: float, existing: dict | None) -> dict:
    extracted: list[dict] = []
    for idx, shape in enumerate(iter_shapes(slide.shapes), start=1):
        box = bbox_norm(shape, slide_w, slide_h)
        st = getattr(shape, "shape_type", None)
        extracted.append(
            {
                "id": shape_name(shape) or f"shape_{idx}",
                "text": shape_text(shape),
                "tag": local_tag(shape),
                "shape_type": str(st or ""),
                "shape_type_int": int(st) if st is not None else None,
                "box": box,
            }
        )

    meaningful = [item for item in extracted if not is_noise(item["box"])]
    line_like = [item for item in meaningful if classify_line_like(item)]
    circle_like = [item for item in meaningful if classify_circle_like(item)]
    card_like = [item for item in meaningful if classify_card_like(item)]

    major = [item for item in meaningful if item["box"]["area"] >= 0.004]
    if not major:
        major = meaningful

    if major:
        min_x = min(item["box"]["x"] for item in major)
        min_y = min(item["box"]["y"] for item in major)
        max_x = max(item["box"]["x"] + item["box"]["w"] for item in major)
        max_y = max(item["box"]["y"] + item["box"]["h"] for item in major)
        main_group = {
            "center_x": round((min_x + max_x) / 2, 4),
            "expected_center_x": float(((existing or {}).get("main_group") or {}).get("expected_center_x", 0.5)),
            "tolerance": float(((existing or {}).get("main_group") or {}).get("tolerance", 0.06)),
        }
        occupancy = {
            "ratio": round((max_x - min_x) * (max_y - min_y), 4),
            "min": float(((existing or {}).get("occupancy") or {}).get("min", 0.24)),
            "max": float(((existing or {}).get("occupancy") or {}).get("max", 0.68)),
        }
    else:
        main_group = {"center_x": 0.5, "expected_center_x": 0.5, "tolerance": 0.06}
        occupancy = {"ratio": 0.0, "min": 0.24, "max": 0.68}

    alignment_groups = list((existing or {}).get("alignment_groups", []))
    # Card columns
    for cluster in cluster_by_coordinate(
        [{"id": item["id"], "x": item["box"]["x"]} for item in card_like],
        "x",
        tolerance=0.03,
        min_size=3,
    ):
        alignment_groups.append(
            {
                "axis": "x",
                "label": f"{cluster[0]['id']}_card_column",
                "tolerance": 0.03,
                "coordinates": [round(item["x"], 4) for item in cluster],
            }
        )
    # Timeline nodes around a vertical line
    vertical_axes = [item for item in line_like if item["box"]["h"] > item["box"]["w"] * 3 and item["box"]["h"] > 0.18]
    for axis in vertical_axes:
        line_cx = axis["box"]["cx"]
        related = [
            item for item in circle_like
            if abs(item["box"]["cx"] - line_cx) <= 0.05
            and axis["box"]["y"] - 0.03 <= item["box"]["cy"] <= axis["box"]["y"] + axis["box"]["h"] + 0.03
        ]
        if len(related) >= 2:
            alignment_groups.append(
                {
                    "axis": "x",
                    "label": f"{axis['id']}_timeline_nodes",
                    "tolerance": 0.025,
                    "coordinates": [round(line_cx, 4), *[round(item["box"]["cx"], 4) for item in related]],
                }
            )

    # Cards
    cards = []
    for item in card_like:
        cards.append(
            {
                "id": item["id"],
                "label": item["text"][:40] or item["id"],
                "height": round(item["box"]["h"], 4),
                "min_height": 0.06,
                "max_height": 0.32,
            }
        )

    # Connectors
    connectors = []
    target_shapes = [item for item in meaningful if item not in line_like]
    for item in line_like:
        # Skip pure timeline axes handled by alignment rules
        if any(group["label"].startswith(item["id"]) for group in alignment_groups):
            continue
        start, end = infer_line_endpoints(item["box"])
        connectors.append(
            {
                "label": item["id"],
                "tolerance": 0.035,
                "start": start,
                "end": end,
                "from_anchor": nearest_anchor(start, target_shapes),
                "to_anchor": nearest_anchor(end, target_shapes),
            }
        )

    entry = {
        "page_id": page_id,
        "role": role,
        "archetype": (existing or {}).get("archetype") or ROLE_TO_ARCHETYPE.get(role, "content_board"),
        "main_group": main_group,
        "occupancy": occupancy,
        "alignment_groups": alignment_groups,
        "connectors": connectors,
        "cards": cards,
        "source": {"type": "pptx", "slide_name": getattr(slide, "name", "")},
    }
    return entry


def merge_entries(existing: dict | None, new_entries: list[dict]) -> dict:
    lookup = {item.get("page_id"): item for item in (existing or {}).get("pages", []) if isinstance(item, dict) and item.get("page_id")}
    for entry in new_entries:
        lookup[entry["page_id"]] = {**lookup.get(entry["page_id"], {}), **entry}
    return {"pages": [lookup[key] for key in sorted(lookup.keys())]}


def build_manifest_from_pptx(deck_path: Path, state: dict, existing: dict | None = None) -> dict:
    prs = Presentation(str(deck_path))
    slide_w = float(prs.slide_width)
    slide_h = float(prs.slide_height)
    role_lookup = {page.get("page_id"): page.get("role", "unassigned") for page in state.get("pages", [])}
    existing_lookup = {item.get("page_id"): item for item in (existing or {}).get("pages", []) if isinstance(item, dict)}
    entries: list[dict] = []
    for idx, slide in enumerate(prs.slides, start=1):
        page_id = f"slide_{idx:02d}"
        role = role_lookup.get(page_id, "unassigned")
        entries.append(build_page_entry(slide, page_id, role, slide_w, slide_h, existing_lookup.get(page_id)))
    return merge_entries(existing, entries)


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract layout geometry from a built PPTX and update layout_manifest.json.")
    parser.add_argument("--project-dir")
    parser.add_argument("--deck-path")
    parser.add_argument("--state")
    parser.add_argument("--manifest")
    args = parser.parse_args()

    project_dir = Path(args.project_dir).expanduser().resolve() if args.project_dir else None
    deck_path = Path(args.deck_path).expanduser().resolve() if args.deck_path else None
    state_path = Path(args.state).expanduser().resolve() if args.state else (project_dir / "slide_state.json" if project_dir else None)
    manifest_path = Path(args.manifest).expanduser().resolve() if args.manifest else (project_dir / "layout_manifest.json" if project_dir else None)
    if not deck_path or not deck_path.exists():
        raise SystemExit("[ERROR] 需要提供存在的 --deck-path。")
    if not state_path or not state_path.exists():
        raise SystemExit("[ERROR] 需要提供存在的 --state 或 --project-dir。")
    if not manifest_path:
        raise SystemExit("[ERROR] 需要提供 --manifest 或 --project-dir。")

    state = load_json(state_path)
    existing = load_json(manifest_path) if manifest_path.exists() else {"pages": []}
    manifest = build_manifest_from_pptx(deck_path, state, existing)
    save_json(manifest_path, manifest)
    print(f"[OK] extracted layout manifest from {deck_path} -> {manifest_path}")


if __name__ == "__main__":
    main()
