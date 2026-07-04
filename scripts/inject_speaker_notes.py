#!/usr/bin/env python3
"""Inject speaker notes from deck_clean_pages.md into existing PPTX files.

Also outputs speaker_notes.json for HTML deck integration.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from page_parser import extract_speaker_notes, extract_speaker_scripts
from validate_external_language_contract import first_forbidden_term, load_forbidden_terms

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def inject_into_pptx(pptx_path: Path, notes: dict[int, str], output_path: Path | None = None) -> int:
    if not HAS_PPTX:
        print("[SKIP] python-pptx not installed, cannot inject speaker notes.")
        return 0

    prs = Presentation(str(pptx_path))
    injected = 0
    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in notes:
            continue
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes[idx]
        injected += 1

    target = output_path or pptx_path
    prs.save(str(target))
    return injected


def write_notes_json(notes: dict[int, str], output_path: Path) -> None:
    payload = {f"slide_{k:02d}": v for k, v in sorted(notes.items())}
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def validate_notes(
    notes: dict[int, str],
    forbidden_terms: list[str],
    *,
    project_dir: Path,
) -> None:
    for page_no, note in sorted(notes.items()):
        term = first_forbidden_term(note, forbidden_terms)
        if term:
            payload = {
                "file": "deck_clean_pages.md",
                "page_id": f"slide_{page_no:02d}",
                "field": "speaker_script",
                "forbidden_term": term,
                "next_command": (
                    "python3 scripts/run_deck_pipeline.py handoff "
                    f"--project-dir {project_dir} --role external-expression"
                ),
            }
            raise SystemExit("[ERROR] speaker script contains internal language:\n" + json.dumps(payload, ensure_ascii=False, indent=2))


def main() -> None:
    parser = argparse.ArgumentParser(description="Inject speaker notes into PPTX and export speaker_notes.json.")
    parser.add_argument("--project-dir", required=True)
    parser.add_argument("--clean-pages")
    parser.add_argument("--pptx-path", help="PPTX file to inject notes into")
    parser.add_argument("--pptx-output", help="Output PPTX path (default: overwrite input)")
    parser.add_argument("--json-output", help="Output speaker_notes.json path")
    parser.add_argument("--language-contract", help="Path to audience_language_contract.json")
    parser.add_argument("--allow-missing-language-contract", action="store_true", help="Use default forbidden terms if no language contract exists")
    parser.add_argument("--legacy-speaker-notes", action="store_true", help="Allow legacy `> 演讲备注:` as a speaker script source")
    args = parser.parse_args()

    project_dir = Path(args.project_dir).expanduser().resolve()
    clean_pages = Path(args.clean_pages or project_dir / "deck_clean_pages.md").expanduser().resolve()
    language_contract = Path(args.language_contract or project_dir / "audience_language_contract.json").expanduser().resolve()

    if not clean_pages.exists():
        raise SystemExit(f"[ERROR] clean pages not found: {clean_pages}")
    if not language_contract.exists() and not args.allow_missing_language_contract:
        raise SystemExit(f"[ERROR] language contract not found: {language_contract}")

    clean_pages_text = clean_pages.read_text(encoding="utf-8")
    legacy_notes = extract_speaker_notes(clean_pages_text)
    if legacy_notes and not args.legacy_speaker_notes:
        raise SystemExit("[ERROR] legacy speaker note format detected. Use `> 讲者话术:` or pass --legacy-speaker-notes.")

    notes = extract_speaker_scripts(clean_pages_text, allow_legacy=args.legacy_speaker_notes)
    if not notes:
        print("[OK] no speaker scripts found in clean pages")
        return
    forbidden_terms = load_forbidden_terms(language_contract if language_contract.exists() else None)
    validate_notes(notes, forbidden_terms, project_dir=project_dir)

    # Export JSON (always)
    json_output = Path(args.json_output or project_dir / "speaker_notes.json").expanduser().resolve()
    write_notes_json(notes, json_output)
    print(f"[OK] wrote {json_output} ({len(notes)} note(s))")

    # Inject into PPTX (if available)
    if args.pptx_path:
        pptx_path = Path(args.pptx_path).expanduser().resolve()
        pptx_output = Path(args.pptx_output).expanduser().resolve() if args.pptx_output else None
        injected = inject_into_pptx(pptx_path, notes, pptx_output)
        print(f"[OK] injected {injected} note(s) into {pptx_output or pptx_path}")


if __name__ == "__main__":
    main()
